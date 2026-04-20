"""Inspect current state of user-edited deck so we append correctly."""
import sys
from pptx import Presentation
sys.stdout.reconfigure(encoding="utf-8")

p = Presentation("Providence - Model Evaluation.pptx")
print(f"TOTAL: {len(p.slides)}")
for i, s in enumerate(p.slides, 1):
    title = ""
    if s.shapes.title is not None and s.shapes.title.has_text_frame:
        title = s.shapes.title.text_frame.text.strip().replace("\n", " / ")
    if not title:
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                title = sh.text_frame.text.strip().split("\n")[0]
                break
    print(f"[{i:02d}] layout='{s.slide_layout.name}' :: {title[:85]}")
