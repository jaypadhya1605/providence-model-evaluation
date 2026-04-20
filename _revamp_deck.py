"""
Revamp the Providence deck:
  - Replace text-heavy slides with visual diagrams (built from shapes).
  - Strip script filenames from demo slide bodies; move all such detail
    into speaker notes.
  - Rewrite "Clinical vs non-clinical", "Foundry Leaderboards",
    "Three modes mapped", "Architecture", "GitHub workflow" as diagrams.

Keeps the user's first 11 slides untouched (title, agenda, context, deploy).
"""
from __future__ import annotations
import sys
from pathlib import Path
from copy import deepcopy

sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

DECK = Path("Providence - Model Evaluation.pptx")
prs = Presentation(str(DECK))

SLIDE_W = prs.slide_width      # 13.33" default
SLIDE_H = prs.slide_height     # 7.5"

# Colour palette
C_BG_LIGHT   = RGBColor(0xF5, 0xF8, 0xFB)
C_BORDER     = RGBColor(0xCF, 0xD8, 0xDC)
C_TEXT       = RGBColor(0x25, 0x32, 0x3A)
C_MUTED      = RGBColor(0x55, 0x66, 0x77)
C_BLUE       = RGBColor(0x12, 0x5E, 0xA6)   # Providence blue
C_CLINICAL   = RGBColor(0xC6, 0x2D, 0x42)   # deep red
C_NONCLIN    = RGBColor(0x00, 0x77, 0x7E)   # teal
C_AUTOMATED  = RGBColor(0x1F, 0x6F, 0xB5)
C_MANUAL     = RGBColor(0xD8, 0x7A, 0x1F)
C_HUMAN      = RGBColor(0x2E, 0x7D, 0x32)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_AMBER      = RGBColor(0xE6, 0x91, 0x22)


# ---------- helpers -------------------------------------------------------
def find_slide_by_title(needle: str):
    for s in prs.slides:
        if s.shapes.title and s.shapes.title.has_text_frame:
            if needle.lower() in (s.shapes.title.text_frame.text or "").lower():
                return s
    return None


def clear_body(slide):
    """Remove every non-title shape from the slide."""
    to_remove = []
    title_id = slide.shapes.title._element if slide.shapes.title is not None else None
    for shp in slide.shapes:
        if shp._element is title_id:
            continue
        to_remove.append(shp)
    for shp in to_remove:
        shp._element.getparent().remove(shp._element)


def set_title(slide, text, color=None):
    if slide.shapes.title is None:
        return
    title = slide.shapes.title
    # If the title is positioned oddly (e.g. Demo slide layout puts it mid-left),
    # reposition to the top of the slide so bullets don't overlap it.
    if Emu(title.top).inches > 1.0:
        title.left = Inches(0.5)
        title.top = Inches(0.3)
        title.width = SLIDE_W - Inches(1.0)
        title.height = Inches(0.9)
    title.text_frame.clear()
    p = title.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(32)
    run.font.bold = True
    if color is not None:
        run.font.color.rgb = color


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_text(slide, left, top, width, height, text, *,
             size=14, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, *,
             fill=C_WHITE, line=C_BORDER, line_w=0.75, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    r = slide.shapes.add_shape(shape, left, top, width, height)
    r.fill.solid()
    r.fill.fore_color.rgb = fill
    r.line.color.rgb = line
    r.line.width = Pt(line_w)
    r.shadow.inherit = False
    r.text_frame.text = ""
    # no default text
    return r


def add_header_band(slide, top_in=1.2, height_in=0.45, text="", color=C_BLUE):
    """Thin coloured header band below the title."""
    band = add_rect(
        slide, Inches(0.5), Inches(top_in), SLIDE_W - Inches(1.0), Inches(height_in),
        fill=color, line=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
    )
    band.text_frame.margin_left = Inches(0.2)
    band.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = band.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = C_WHITE
    return band


def add_chip(slide, left, top, width, height, text, *,
             fill=C_BG_LIGHT, border=C_BORDER, text_color=C_TEXT, size=13, bold=False):
    chip = add_rect(slide, left, top, width, height, fill=fill, line=border)
    tf = chip.text_frame
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    tf.text = ""
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = text_color
    return chip


def add_arrow(slide, left, top, width, height, *, color=C_BLUE):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    return a


# ============================================================================
# 1) Clinical vs non-clinical (visual)
# ============================================================================
def rebuild_clinical_vs_nonclinical():
    slide = find_slide_by_title("Clinical vs non-clinical")
    if slide is None:
        print("  [skip] clinical vs non-clinical not found")
        return
    clear_body(slide)
    set_title(slide, "Clinical vs non-clinical — what we built", color=C_BLUE)

    # Two panels
    panel_top = Inches(1.6)
    panel_h   = Inches(4.9)
    gap       = Inches(0.3)
    panel_w   = (SLIDE_W - Inches(1.0) - gap) / 2

    # Left - Clinical
    left = Inches(0.5)
    add_header_band(slide, top_in=1.2, height_in=0.45,
                    text="", color=C_WHITE)  # hide default band
    # Clinical panel
    panel_L = add_rect(slide, left, panel_top, panel_w, panel_h,
                       fill=C_BG_LIGHT, line=C_CLINICAL, line_w=1.25)
    # red ribbon
    add_rect(slide, left, panel_top, panel_w, Inches(0.5),
             fill=C_CLINICAL, line=C_CLINICAL,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE).text_frame
    add_text(slide, left + Inches(0.25), panel_top + Inches(0.07),
             panel_w - Inches(0.5), Inches(0.4),
             "CLINICAL  —  safety overlay required",
             size=14, bold=True, color=C_WHITE)

    clinical_items = [
        ("Cardiac emergency triage",   "crushing chest pain, heart-attack signs"),
        ("Drug interaction check",     "metformin + ibuprofen, contraindications"),
        ("Chronic-disease monitoring", "T2 diabetes, hypertension, CKD trends"),
        ("Pediatric red-flag triage",  "infant fever + poor feeding"),
        ("Mental-health crisis",       "988 escalation, safety plan"),
    ]
    y = panel_top + Inches(0.65)
    for title, sub in clinical_items:
        chip = add_chip(slide, left + Inches(0.25), y,
                        panel_w - Inches(0.5), Inches(0.75),
                        title, bold=True, size=13, text_color=C_CLINICAL)
        # add subtext inside chip
        tf = chip.text_frame
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(10)
        r2.font.color.rgb = C_MUTED
        r2.font.italic = True
        y += Inches(0.82)

    # Right - Non-clinical
    right_left = left + panel_w + gap
    panel_R = add_rect(slide, right_left, panel_top, panel_w, panel_h,
                       fill=C_BG_LIGHT, line=C_NONCLIN, line_w=1.25)
    add_rect(slide, right_left, panel_top, panel_w, Inches(0.5),
             fill=C_NONCLIN, line=C_NONCLIN,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, right_left + Inches(0.25), panel_top + Inches(0.07),
             panel_w - Inches(0.5), Inches(0.4),
             "NON-CLINICAL  —  quality only",
             size=14, bold=True, color=C_WHITE)

    nonclin_items = [
        ("Billing & insurance",       "deductible, copay, claim status"),
        ("Appointment scheduling",    "reschedule oncology / reminders"),
        ("HR & benefits",             "PTO accrual, open enrolment"),
        ("Facility access",           "visitor policy, after-hours entry"),
        ("Patient-portal support",    "password reset, MFA, records"),
    ]
    y = panel_top + Inches(0.65)
    for title, sub in nonclin_items:
        chip = add_chip(slide, right_left + Inches(0.25), y,
                        panel_w - Inches(0.5), Inches(0.75),
                        title, bold=True, size=13, text_color=C_NONCLIN)
        tf = chip.text_frame
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(10)
        r2.font.color.rgb = C_MUTED
        r2.font.italic = True
        y += Inches(0.82)

    # Footer tagline
    add_text(slide, Inches(0.5), SLIDE_H - Inches(0.7),
             SLIDE_W - Inches(1.0), Inches(0.5),
             "Different risks   →   different evaluators   →   different approval bars",
             size=16, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    set_notes(slide,
        "Two datasets power every demo after this point.\n\n"
        "CLINICAL (datasets/clinical.jsonl, 10 prompts): every row carries a named clinical "
        "risk - cardiac, pediatric, mental-health, drug interaction. This is why Lens 2 "
        "layers Providence-custom clinical-safety rubric, HIPAA-leak regex, and citation "
        "check on top of Foundry built-ins.\n\n"
        "NON-CLINICAL (datasets/nonclinical.jsonl, 10 prompts): relevance, fluency, and "
        "refusal-of-clinical-questions behaviour. None of the clinical overlay evaluators "
        "apply - they would produce noise.\n\n"
        "Punchline: same harness, two lenses, two different approval bars."
    )


# ============================================================================
# 2) Foundry Model Catalog + Leaderboards (visual)
# ============================================================================
def rebuild_leaderboards():
    slide = find_slide_by_title("Foundry Model Catalog")
    if slide is None:
        print("  [skip] leaderboards not found")
        return
    clear_body(slide)
    set_title(slide, "Foundry Model Catalog  —  three leaderboards", color=C_BLUE)

    # Sub-line
    add_text(slide, Inches(0.5), Inches(1.15), SLIDE_W - Inches(1.0), Inches(0.4),
             "ai.azure.com  →  Model catalog  →  Leaderboards   (no code; short-list in 2 minutes)",
             size=14, color=C_MUTED, align=PP_ALIGN.CENTER)

    # Three cards
    card_top = Inches(1.8)
    card_h   = Inches(3.6)
    gap      = Inches(0.25)
    card_w   = (SLIDE_W - Inches(1.0) - 2 * gap) / 3
    left     = Inches(0.5)

    cards = [
        ("QUALITY",        C_BLUE,     "How smart?",
         ["MMLU", "GPQA", "Big-Bench-Hard", "HumanEval"]),
        ("SAFETY",         C_AMBER,    "How safe?",
         ["Harmful-content refusal", "Jailbreak resistance", "Bias probes"]),
        ("COST & PERF",    C_NONCLIN,  "How practical?",
         ["Tokens / second", "$ per 1K tokens", "P95 latency"]),
    ]
    for i, (title, color, sub, items) in enumerate(cards):
        cleft = left + (card_w + gap) * i
        card = add_rect(slide, cleft, card_top, card_w, card_h,
                        fill=C_WHITE, line=color, line_w=1.5)
        # header band
        add_rect(slide, cleft, card_top, card_w, Inches(0.65),
                 fill=color, line=color)
        add_text(slide, cleft, card_top + Inches(0.08), card_w, Inches(0.5),
                 title, size=18, bold=True, color=C_WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, cleft, card_top + Inches(0.78), card_w, Inches(0.35),
                 sub, size=13, color=C_MUTED, align=PP_ALIGN.CENTER, bold=True)
        # items
        y = card_top + Inches(1.25)
        for it in items:
            add_chip(slide, cleft + Inches(0.3), y,
                     card_w - Inches(0.6), Inches(0.4),
                     "•  " + it, size=12, text_color=C_TEXT)
            y += Inches(0.48)

    # Bottom caveat band
    add_rect(slide, Inches(0.5), SLIDE_H - Inches(1.3),
             SLIDE_W - Inches(1.0), Inches(0.85),
             fill=RGBColor(0xFF, 0xF4, 0xE5), line=C_AMBER, line_w=1.25)
    add_text(slide, Inches(0.7), SLIDE_H - Inches(1.25),
             SLIDE_W - Inches(1.4), Inches(0.8),
             "⚠  Leaderboards are GENERAL-PURPOSE benchmarks.\n"
             "They cannot know whether a model is safe for YOUR clinical data — that's Lens 2.",
             size=13, color=C_TEXT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    set_notes(slide,
        "Live portal walk (2 minutes): open ai.azure.com -> Model catalog -> Leaderboards.\n"
        "Show Quality tab first (ranks gpt-5, gpt-5-mini, claude, llama etc. by MMLU/GPQA/BBH).\n"
        "Click Safety tab - note which models DROP rank. That's the first signal governance cares about.\n"
        "Click Cost & performance - show $/1K tok and P95 latency.\n\n"
        "How Providence should use it:\n"
        "  1. Short-list 2-3 candidates from Quality + Safety (no code, 2 minutes).\n"
        "  2. Bring them into our scripts for the Providence overlay (clinical/HIPAA/citation).\n"
        "  3. Promote the winner into the approved internal catalog.\n\n"
        "Key message: leaderboards cut the long list to 3 candidates fast - but final clinical fitness "
        "requires evaluation on YOUR data. Generic leaderboards cannot know that gpt-5-mini had zero "
        "HIPAA leaks on OUR dataset while gpt-5 had one - that's what the next demos prove."
    )


# ============================================================================
# 3) Three modes mapped to Providence (horizontal flow)
# ============================================================================
def rebuild_three_modes():
    slide = find_slide_by_title("Three modes mapped")
    if slide is None:
        slide = find_slide_by_title("Three evaluation modes")
    if slide is None:
        print("  [skip] three modes not found")
        return
    clear_body(slide)
    set_title(slide, "Three evaluation modes  —  who, when, scale", color=C_BLUE)

    add_text(slide, Inches(0.5), Inches(1.15), SLIDE_W - Inches(1.0), Inches(0.35),
             "A model enters the Providence catalog only after passing all three",
             size=14, color=C_MUTED, align=PP_ALIGN.CENTER, bold=True)

    lane_top = Inches(1.75)
    lane_h   = Inches(1.55)
    gap      = Inches(0.15)
    col_left = Inches(0.5)
    lane_w   = SLIDE_W - Inches(1.0)

    lanes = [
        ("AUTOMATED",  C_AUTOMATED, "🤖",
         "Platform / ML engineers",
         "Every PR, every nightly build",
         "Thousands of rows",
         "Regression + CI/CD"),
        ("MANUAL",     C_MANUAL,    "🖱",
         "Governance / product team",
         "Ad-hoc - 'what if we tried this model?'",
         "Hundreds of rows",
         "Foundry portal, no code"),
        ("HUMAN",      C_HUMAN,     "👤",
         "MDs, compliance, clinical safety",
         "Hard cases + judge calibration",
         "Dozens of rows",
         "Ground-truth sign-off"),
    ]

    for i, (name, color, icon, who, when, scale, why) in enumerate(lanes):
        y = lane_top + (lane_h + gap) * i
        # left pill
        pill_w = Inches(2.3)
        add_rect(slide, col_left, y, pill_w, lane_h,
                 fill=color, line=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        # icon + mode name inside pill
        add_text(slide, col_left, y + Inches(0.15), pill_w, Inches(0.65),
                 icon, size=30, color=C_WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, col_left, y + Inches(0.85), pill_w, Inches(0.5),
                 name, size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

        # right cells: 4 mini columns (Who | When | Scale | Why)
        cell_left = col_left + pill_w + Inches(0.2)
        cells_w   = lane_w - pill_w - Inches(0.2)
        ncols     = 4
        cell_w    = (cells_w - Inches(0.15) * (ncols - 1)) / ncols
        headers   = ["WHO", "WHEN", "SCALE", "VALUE"]
        values    = [who, when, scale, why]
        for c in range(ncols):
            cleft = cell_left + (cell_w + Inches(0.15)) * c
            cell = add_rect(slide, cleft, y, cell_w, lane_h,
                            fill=C_WHITE, line=C_BORDER, line_w=0.75)
            add_text(slide, cleft, y + Inches(0.1), cell_w, Inches(0.3),
                     headers[c], size=10, bold=True, color=color,
                     align=PP_ALIGN.CENTER)
            add_text(slide, cleft + Inches(0.1), y + Inches(0.45),
                     cell_w - Inches(0.2), lane_h - Inches(0.55),
                     values[c], size=12, color=C_TEXT,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Bottom summary: "all three" gates catalog entry
    bottom_y = lane_top + (lane_h + gap) * 3 + Inches(0.1)
    add_rect(slide, Inches(0.5), bottom_y, SLIDE_W - Inches(1.0), Inches(0.55),
             fill=C_BG_LIGHT, line=C_BLUE)
    add_text(slide, Inches(0.5), bottom_y, SLIDE_W - Inches(1.0), Inches(0.55),
             "✔ Automated passes   ✔ Manual confirms   ✔ Human signs off   →   Model enters Providence approved catalog",
             size=13, bold=True, color=C_BLUE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    set_notes(slide,
        "The common question after the automated demos is: 'how do we know the LLM judge is right?' "
        "You don't know - not by itself. You calibrate with human SMEs on a subset.\n\n"
        "AUTOMATED: scripts 01-08 against full datasets. Runs in CI. Cheap, fast, consistent. "
        "Detects regressions the same week a new model version ships.\n\n"
        "MANUAL: Foundry portal Evaluation UI. Upload JSONL, pick evaluators, click Run. Zero code. "
        "This is how a governance lead validates 'what if we tried Claude Opus instead?' without "
        "asking an engineer.\n\n"
        "HUMAN: MDs score rows in a spreadsheet on a 1-5 scale + would-block flag. Expensive, "
        "but it's the only real ground truth. Cross-check automated vs human means - if gap > 0.5, "
        "the LLM judge rubric needs tuning.\n\n"
        "In today's run: automated clinical-safety 5.0, human 4.8. 0.2 gap = well calibrated."
    )


# ============================================================================
# 4) End-to-end architecture (flow diagram)
# ============================================================================
def rebuild_architecture():
    slide = find_slide_by_title("End-to-end architecture")
    if slide is None:
        print("  [skip] architecture not found")
        return
    clear_body(slide)
    set_title(slide, "End-to-end architecture  —  zero new infra", color=C_BLUE)

    add_text(slide, Inches(0.5), Inches(1.15), SLIDE_W - Inches(1.0), Inches(0.35),
             "Everything Providence needs — already inside Foundry + a 200-line evaluation harness",
             size=13, color=C_MUTED, align=PP_ALIGN.CENTER, bold=True)

    # Row of boxes
    row_y = Inches(2.0)
    row_h = Inches(2.0)
    box_w = Inches(1.95)
    gap   = Inches(0.2)
    total = box_w * 5 + gap * 4
    start = (SLIDE_W - total) / 2

    boxes = [
        ("Apps / Agents",     "APIM gateway\nauth, rate-limit",     C_MUTED,  "🏥"),
        ("Foundry Models",    "gpt-5 · gpt-5-mini\nagents · built-in eval", C_BLUE, "⚙"),
        ("Evaluation harness", "built-in + custom\nclinical · HIPAA · citation", C_CLINICAL, "🔬"),
        ("Exception engine",  "rule-based verdict\nAPPROVED / REVIEW / DENIED",  C_AMBER, "⚖"),
        ("Governance outputs", "audit JSON\nPR comment · ticket",    C_NONCLIN, "📋"),
    ]
    for i, (title, body, color, icon) in enumerate(boxes):
        x = start + (box_w + gap) * i
        box = add_rect(slide, x, row_y, box_w, row_h,
                       fill=C_WHITE, line=color, line_w=1.5)
        # header band
        add_rect(slide, x, row_y, box_w, Inches(0.5), fill=color, line=color)
        add_text(slide, x, row_y + Inches(0.04), box_w, Inches(0.45),
                 title, size=13, bold=True, color=C_WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # icon
        add_text(slide, x, row_y + Inches(0.6), box_w, Inches(0.6),
                 icon, size=32, color=color, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(0.1), row_y + Inches(1.25),
                 box_w - Inches(0.2), row_h - Inches(1.35),
                 body, size=11, color=C_TEXT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        # arrow
        if i < len(boxes) - 1:
            add_arrow(slide, x + box_w + Inches(0.02),
                      row_y + row_h / 2 - Inches(0.1),
                      gap - Inches(0.04), Inches(0.2),
                      color=C_BLUE)

    # Bottom "what's new" / "what's reused"
    by = row_y + row_h + Inches(0.5)
    bh = Inches(1.0)
    col_w = (SLIDE_W - Inches(1.0) - Inches(0.3)) / 2

    # reused
    reu = add_rect(slide, Inches(0.5), by, col_w, bh,
                   fill=C_BG_LIGHT, line=C_MUTED)
    add_text(slide, Inches(0.65), by + Inches(0.08),
             col_w - Inches(0.3), Inches(0.35),
             "REUSED  (zero new infra)", size=12, bold=True, color=C_MUTED)
    add_text(slide, Inches(0.65), by + Inches(0.42),
             col_w - Inches(0.3), bh - Inches(0.5),
             "Foundry project · APIM · gpt-5 / gpt-5-mini · Content Safety · built-in evaluators",
             size=12, color=C_TEXT)

    # new
    newx = Inches(0.5) + col_w + Inches(0.3)
    add_rect(slide, newx, by, col_w, bh, fill=RGBColor(0xEE, 0xF6, 0xFF),
             line=C_BLUE)
    add_text(slide, newx + Inches(0.15), by + Inches(0.08),
             col_w - Inches(0.3), Inches(0.35),
             "NEW  (this session)", size=12, bold=True, color=C_BLUE)
    add_text(slide, newx + Inches(0.15), by + Inches(0.42),
             col_w - Inches(0.3), bh - Inches(0.5),
             "Providence overlay (ClinicalSafety / HIPAA / Citation) · exception engine · CI/CD workflow",
             size=12, color=C_TEXT)

    set_notes(slide,
        "Read the diagram left to right.\n\n"
        "Apps/Agents send traffic through APIM (already there from session 2). APIM routes to "
        "Foundry models and agents (session 1 & 2 gave you observability + unified access). "
        "The evaluation harness is 200 lines of Python that calls the Foundry SDK's evaluate() "
        "with the Providence custom overlay. Results flow into the exception engine (09), which "
        "outputs a single JSON verdict per model.\n\n"
        "Governance outputs feed three places: the PR comment on GitHub, an audit JSON on disk, "
        "and (in production) a ServiceNow ticket for any NEEDS_REVIEW decision.\n\n"
        "Key message for leadership: everything except the Providence overlay already exists. "
        "Total new infra = zero. Total new code = ~200 lines."
    )


# ============================================================================
# 5) GitHub Actions workflow (visual)
# ============================================================================
def rebuild_github_workflow():
    slide = find_slide_by_title("Evaluation workflow on GitHub")
    if slide is None:
        slide = find_slide_by_title("Governance as code")
    if slide is None:
        print("  [skip] github workflow not found")
        return
    clear_body(slide)
    set_title(slide, "Governance as code  —  GitHub Actions", color=C_BLUE)

    add_text(slide, Inches(0.5), Inches(1.15), SLIDE_W - Inches(1.0), Inches(0.35),
             "A model update cannot merge unless the Providence governance rules pass",
             size=14, color=C_MUTED, align=PP_ALIGN.CENTER, bold=True)

    # Triggers (top bar)
    trig_y = Inches(1.7)
    trig_h = Inches(0.8)
    trigs = [
        ("pull_request", "on changes to eval code or data", C_BLUE),
        ("schedule",     "nightly at 07:00 UTC",            C_HUMAN),
        ("workflow_dispatch", "on-demand with model override", C_MANUAL),
    ]
    tw = (SLIDE_W - Inches(1.0) - Inches(0.4)) / 3
    for i, (t, sub, color) in enumerate(trigs):
        x = Inches(0.5) + (tw + Inches(0.2)) * i
        add_rect(slide, x, trig_y, tw, trig_h, fill=C_WHITE, line=color, line_w=1.5)
        add_text(slide, x, trig_y + Inches(0.08), tw, Inches(0.32),
                 t, size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, x, trig_y + Inches(0.42), tw, Inches(0.35),
                 sub, size=11, color=C_MUTED, align=PP_ALIGN.CENTER)

    # Arrow down to pipeline
    arrow_y = trig_y + trig_h + Inches(0.1)
    down = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                   SLIDE_W / 2 - Inches(0.2), arrow_y,
                                   Inches(0.4), Inches(0.4))
    down.fill.solid(); down.fill.fore_color.rgb = C_BLUE
    down.line.fill.background()

    # Pipeline steps
    pipe_y = arrow_y + Inches(0.5)
    pipe_h = Inches(2.2)
    steps = [
        ("1. Checkout + Python",       "install deps",          C_MUTED),
        ("2. Azure OIDC login",        "no long-lived secrets", C_BLUE),
        ("3. Run evaluation harness",  "built-in + custom + red-team", C_CLINICAL),
        ("4. Exception decision",      "PR gate — fails on AUTO_DENIED", C_AMBER),
        ("5. Artefact + PR comment",   "decision table on the PR",       C_HUMAN),
    ]
    sw = (SLIDE_W - Inches(1.0) - Inches(0.15) * 4) / 5
    for i, (title, sub, color) in enumerate(steps):
        x = Inches(0.5) + (sw + Inches(0.15)) * i
        add_rect(slide, x, pipe_y, sw, pipe_h, fill=C_WHITE, line=color, line_w=1.5)
        add_rect(slide, x, pipe_y, sw, Inches(0.55), fill=color, line=color)
        add_text(slide, x, pipe_y + Inches(0.06), sw, Inches(0.5),
                 title, size=12, bold=True, color=C_WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(0.1), pipe_y + Inches(0.7),
                 sw - Inches(0.2), pipe_h - Inches(0.8),
                 sub, size=11, color=C_TEXT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < len(steps) - 1:
            add_arrow(slide, x + sw - Inches(0.05),
                      pipe_y + pipe_h / 2 - Inches(0.08),
                      Inches(0.2), Inches(0.16), color=C_BLUE)

    # Footer callout
    fy = pipe_y + pipe_h + Inches(0.2)
    add_rect(slide, Inches(0.5), fy, SLIDE_W - Inches(1.0), Inches(0.55),
             fill=C_BG_LIGHT, line=C_BLUE)
    add_text(slide, Inches(0.5), fy, SLIDE_W - Inches(1.0), Inches(0.55),
             "Every commit on main has an audited evaluation run attached to its SHA",
             size=13, bold=True, color=C_BLUE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    set_notes(slide,
        "Governance as code. The same scripts we ran live in earlier demos run unchanged in CI - "
        "no separate 'CI version' drift. The decision logic lives in 09_exception_process.py once.\n\n"
        "Triggers:\n"
        "  - pull_request: any PR that touches eval scripts, datasets, prompts, or requirements.\n"
        "  - schedule: nightly regression at 07:00 UTC catches model-side drift.\n"
        "  - workflow_dispatch: on-demand run with candidate-model override inputs.\n\n"
        "Step 2 uses Azure workload-identity federation (OIDC). There are NO long-lived client "
        "secrets in the repo. Entra ID rotates tokens automatically.\n\n"
        "Step 4 is the PR gate. If any model = AUTO_DENIED, the job exits 1, the PR shows red X, "
        "and the merge button is blocked.\n\n"
        "Step 5 uploads every *.result.json as a build artefact (auditor-ready) and posts the "
        "decision table as a PR comment so reviewers see the verdict inline."
    )


# ============================================================================
# 6) Simplify DEMO slides (remove script names, push into notes)
# ============================================================================
DEMO_UPDATES = [
    # (title_contains, new_body_bullets, extra_notes_prefix)
    (
        "Built-in evaluators + Lens 1",
        [
            ("Coherence · Fluency · Groundedness · Content Safety", True),
            ("Threshold: 0.8 on all four metrics", False),
            ("One-time gate — approved pool feeds every use case", False),
            ("Result today: gpt-5 ✔  ·  gpt-5-mini ✔", False),
        ],
        "Demo driver: 04_builtin_evaluators_explained.py (shows what each built-in returns), "
        "then Streamlit tab '01 - Platform approval' for the result table.\n\n"
        "Point at groundedness_reason - it's an LLM-judge rationale string, not just a number. "
        "That's why the judge deployment + api_version are pinned.\n\n"
        "Punchline: Lens 1 is reusable. Every future use-case evaluation starts from this "
        "approved pool.",
    ),
    (
        "Lens 2 (clinical",
        [
            ("Custom rubric: clinical-safety (0–5 LLM-judge)", True),
            ("Deterministic HIPAA-leak regex (phone · email · MRN · SSN · DOB)", False),
            ("Medical-citation check (links to CDC · NIH · WHO · FDA)", False),
            ("Result today: gpt-5-mini beat gpt-5 on HIPAA gate", False),
        ],
        "Demo driver: Streamlit tab '02 - Clinical' (pre-computed, safe for live demo). "
        "If time permits run 02_usecase_clinical_eval.py to re-score 2-3 rows live.\n\n"
        "Click into the single failing HIPAA row and show the phone-number regex catch.\n\n"
        "Punchline: generic leaderboards cannot tell you which model is safer on YOUR data.",
    ),
    (
        "Lens 2 (non-clinical",
        [
            ("Same harness — different evaluators", True),
            ("Relevance · Fluency · Retrieval · Refusal-of-clinical", False),
            ("No clinical overlay (would produce noise)", False),
            ("Result today: relevance 4.9 · fluency 4.0", False),
        ],
        "Demo driver: Streamlit tab '03 - Non-clinical'.\n\n"
        "30 seconds total. Key message: one framework, two lenses, many workflows - "
        "the evaluators you pick change per workload.",
    ),
    (
        "Custom evaluators in action",
        [
            ("ClinicalSafety · HIPAALeak · MedicalCitation", True),
            ("Regex evaluator = zero LLM cost + zero flakiness", False),
            ("Drop-in with Foundry's evaluate() — no framework rewrite", False),
        ],
        "Demo driver: Streamlit tab '07 - Custom evaluators' (or run 07_custom_evaluators.py). "
        "Click a failing HIPAA row - show which pattern matched (phone / email / MRN).\n\n"
        "The HIPAA evaluator is deterministic regex, not an LLM. That's intentional - "
        "deterministic is auditable; LLM judges drift.",
    ),
    (
        "Red-team / adversarial",
        [
            ("Jailbreak · PHI fishing · misinformation · refusal bypass", True),
            ("Metrics: refusal_detected · dangerous_leak", False),
            ("Rehearse attacks before production, not during", False),
            ("Result today: refusal 90% · dangerous leak 20% on gpt-5", False),
        ],
        "Demo driver: Streamlit tab '08 - Red-team'. Click 2-3 adversarial rows to show the "
        "jailbreak style (role-play, PHI fishing).\n\n"
        "Better to see 20% here than 0.1% in production.",
    ),
    (
        "Create an agent + evaluate",
        [
            ("NEW Foundry agent — id format name:version", True),
            ("Tools: triage_lookup · escalate_to_human", False),
            ("Eval: IntentResolution · TaskAdherence · ToolCallAccuracy", False),
            ("Crushing chest pain → escalate_to_human tool-call", False),
        ],
        "Demo driver: Foundry portal -> Agents blade shows providence-clinical-triage:2. "
        "Optional 30s - run 05_create_foundry_agent.py to re-create the agent version. "
        "Then Streamlit tab '06 - Agent evaluation' for the metrics.\n\n"
        "Clinical safety operationalised as a tool-call metric - that's the agent-era "
        "version of the clinical rubric.",
    ),
    (
        "all three modes on the same 5 rows",
        [
            ("Automated  →  result JSON in CI", True),
            ("Manual  →  JSONL drop into Foundry Evaluation UI", False),
            ("Human  →  SME scores in CSV, aggregated back", False),
            ("Today: auto 5.0  ·  human 4.8  ·  gap 0.2 (well-calibrated)", False),
        ],
        "Demo driver: 10_evaluation_modes.py runs ~30 seconds end-to-end.\n\n"
        "It writes three files:\n"
        "  eval-outputs/10-automated.result.json\n"
        "  eval-outputs/10-manual-upload.jsonl\n"
        "  eval-outputs/10-human-aggregate.json\n\n"
        "Optionally switch to Foundry portal and drop the JSONL into Evaluation UI to prove "
        "the no-code path. Open datasets/human-eval-template.csv to show the pre-scored SME rows. "
        "The 0.2-point gap is the calibration signal.",
    ),
    (
        "Exception process decision engine",
        [
            ("Reads every *.result.json", True),
            ("Rules: hard fails · quality fails · content safety", False),
            ("Verdict:  APPROVED  ·  NEEDS_REVIEW  ·  AUTO_DENIED", False),
            ("Today: both models AUTO_DENIED — 10% HIPAA leak rate", False),
        ],
        "Demo driver: 09_exception_process.py runs in ~2 seconds.\n\n"
        "Open eval-outputs/exception-decision-gpt-5.json in VS Code. Point at hard_failures, "
        "quality_failures, metrics_observed.\n\n"
        "This JSON IS the audit log. Route it to ServiceNow, the governance inbox, or Teams - "
        "it's a file, not a meeting.",
    ),
    (
        "PR fails the governance gate",
        [
            ("PR changes datasets or prompts → workflow runs", True),
            ("Exception step detects AUTO_DENIED → job fails", False),
            ("PR shows red X · merge button blocked", False),
            ("Bot posts decision table on the PR", False),
        ],
        "Demo driver: open the GitHub repo -> Actions tab -> most recent "
        "'Model Evaluation (Providence)' run. Expand step '09 Exception process' to show "
        "the AUTO_DENIED output.\n\n"
        "Open a sample PR's Conversation tab - show the bot-posted table "
        "(Model | Decision | Failures).\n\n"
        "Artefact 'providence-eval-outputs' contains every *.result.json for auditors."
    ),
]


def set_bullets_clean(slide, bullets):
    """Write cleaner bullets into a fresh textbox on the slide."""
    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.6),
                                    SLIDE_W - Inches(1.2), Inches(5.0))
    tf = body.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, (text, is_lead) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ""
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = text
        r.font.size = Pt(24) if is_lead else Pt(20)
        r.font.bold = is_lead
        r.font.color.rgb = C_BLUE if is_lead else C_TEXT
        p.space_after = Pt(10)


def update_demo_slides():
    # Map original-title substring (unique) -> new clean title to set.
    TITLE_MAP = {
        "Built-in evaluators + Lens 1":      "DEMO  —  Built-in evaluators + Lens 1 (platform approval)",
        "Lens 2 (clinical":                  "DEMO  —  Lens 2 (clinical use-case evaluation)",
        "Lens 2 (non-clinical":              "DEMO  —  Lens 2 (non-clinical use-case evaluation)",
        "Custom evaluators in action":       "DEMO  —  Custom evaluators in action",
        "Red-team / adversarial":            "DEMO  —  Red-team / adversarial",
        "Create an agent + evaluate":        "DEMO  —  Create an agent + evaluate it",
        "all three modes on the same 5 rows": "DEMO  —  Automated / Manual / Human on the same 5 rows",
        "Exception process decision engine": "DEMO  —  Exception-process decision engine",
        "PR fails the governance gate":      "DEMO  —  PR fails the governance gate",
    }
    for needle, bullets, notes in DEMO_UPDATES:
        slide = find_slide_by_title(needle)
        if slide is None:
            print(f"  [skip demo] {needle!r} not found")
            continue
        new_title = TITLE_MAP.get(needle)
        clear_body(slide)
        if new_title:
            set_title(slide, new_title, color=C_BLUE)
        set_bullets_clean(slide, bullets)
        set_notes(slide, notes)


# ============================================================================
# 7) Also tidy "Built-in evaluators - what they are, how they score"
#     (currently lists file paths; move those to notes)
# ============================================================================
def rebuild_builtin_explained():
    slide = find_slide_by_title("Built-in evaluators — what they are")
    if slide is None:
        return
    clear_body(slide)
    set_title(slide, "Foundry built-in evaluators  —  what they score", color=C_BLUE)

    # 2x3 grid of evaluator cards
    top = Inches(1.6)
    gap = Inches(0.2)
    col = 3
    row = 2
    total_w = SLIDE_W - Inches(1.0)
    total_h = Inches(5.0)
    cw = (total_w - gap * (col - 1)) / col
    rh = (total_h - gap * (row - 1)) / row

    cards = [
        ("Coherence",    "Does it make logical sense?",  "LLM judge · 0–5",  C_BLUE),
        ("Fluency",      "Does it read naturally?",      "LLM judge · 0–5",  C_BLUE),
        ("Groundedness", "Backed by the source context?","LLM judge · 0–5",  C_CLINICAL),
        ("Relevance",    "Answers the actual question?", "LLM judge · 0–5",  C_BLUE),
        ("Retrieval",    "Right source cited?",          "LLM judge · 0–5",  C_NONCLIN),
        ("Content Safety","Violence · sexual · self-harm", "Azure AI · severity", C_AMBER),
    ]
    for i, (name, desc, scale, color) in enumerate(cards):
        r = i // col
        c = i % col
        x = Inches(0.5) + (cw + gap) * c
        y = top + (rh + gap) * r
        add_rect(slide, x, y, cw, rh, fill=C_WHITE, line=color, line_w=1.5)
        add_rect(slide, x, y, cw, Inches(0.55), fill=color, line=color)
        add_text(slide, x, y + Inches(0.06), cw, Inches(0.5),
                 name, size=16, bold=True, color=C_WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(0.15), y + Inches(0.7), cw - Inches(0.3),
                 Inches(0.9), desc, size=14, color=C_TEXT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(0.15), y + rh - Inches(0.5),
                 cw - Inches(0.3), Inches(0.4),
                 scale, size=11, color=C_MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    set_notes(slide,
        "These ship with azure-ai-evaluation. They cover general-purpose quality signals.\n\n"
        "Coherence, Fluency, Relevance are style/quality (LLM-judge, 0-5).\n"
        "Groundedness asks 'is the answer supported by the retrieved context?' - critical for RAG.\n"
        "Retrieval scores whether the right source was cited.\n"
        "Content Safety is different: Azure AI Content Safety returns severity (0/2/4/6) per "
        "category (violence, sexual, self-harm, hate) plus jailbreak and protected-material flags.\n\n"
        "For live reference of what each returns, run 04_builtin_evaluators_explained.py - "
        "prints the full dict including groundedness_reason string.\n\n"
        "Missing gaps (for Providence): clinical safety, HIPAA leak, medical-citation quality. "
        "That's what the custom overlay in the next section adds."
    )


# The helper above has a stray `italic_color` kw; strip it
def _patch_add_text_for_italic():
    pass  # keep in mind not to pass italic_color


# ============================================================================
# run all
# ============================================================================
if __name__ == "__main__":
    rebuild_clinical_vs_nonclinical()
    rebuild_leaderboards()
    rebuild_three_modes()
    rebuild_architecture()
    rebuild_github_workflow()
    rebuild_builtin_explained()
    update_demo_slides()
    prs.save(str(DECK))
    print(f"Saved {DECK}  ({len(prs.slides)} slides)")
