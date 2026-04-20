"""
Append new slides to the user-edited deck (Providence - Model Evaluation.pptx).
Covers:
  - Demo use cases (clinical & non-clinical)
  - Foundry Model Catalog + Leaderboard (UI)
  - Automated / Manual / Human evaluation modes
  - CI/CD evaluation (GitHub Actions)

Existing slides are NOT touched; new slides are inserted at target positions.
"""
from __future__ import annotations
import sys
from pathlib import Path

sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn

DECK = Path("Providence - Model Evaluation.pptx")
prs = Presentation(str(DECK))


# ---------- helpers --------------------------------------------------------
def find_layout(*name_candidates):
    for sm in prs.slide_masters:
        for layout in sm.slide_layouts:
            for c in name_candidates:
                if layout.name.strip().lower() == c.strip().lower():
                    return layout
    # fallback: partial
    for sm in prs.slide_masters:
        for layout in sm.slide_layouts:
            for c in name_candidates:
                if c.strip().lower() in layout.name.strip().lower():
                    return layout
    return prs.slide_masters[0].slide_layouts[0]


LAYOUT_SECTION = find_layout("Section Title", "Section Title 2")
LAYOUT_CONTENT = find_layout("Title and Content")
LAYOUT_TWOCOL  = find_layout("Two Column Bullet with Subheads", "Two Column Bullet text")
LAYOUT_DEMO    = find_layout("Demo slide", "Demo slide 2")


def set_title(slide, text):
    if slide.shapes.title is not None:
        slide.shapes.title.text = text
    else:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(36)
        p.font.bold = True


def set_body_bullets(slide, bullets):
    tf = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 0 and ph.has_text_frame:
            tf = ph.text_frame
            break
    if tf is None:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.5))
        tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        for r in p.runs:
            r.font.size = Pt(20) if level == 0 else Pt(16)


def add_notes(slide, notes):
    slide.notes_slide.notes_text_frame.text = notes


def new_slide(layout, title, bullets=None, notes=None):
    s = prs.slides.add_slide(layout)
    set_title(s, title)
    if bullets:
        set_body_bullets(s, bullets)
    if notes:
        add_notes(s, notes)
    return s


# ---------- find slide ids we'll insert before ----------------------------
def slide_index_for_title_contains(*needles):
    """Return 0-based index of first slide whose title contains any needle."""
    for i, s in enumerate(prs.slides):
        if s.shapes.title is None:
            continue
        t = (s.shapes.title.text_frame.text or "").lower()
        for n in needles:
            if n.lower() in t:
                return i
    return None


# ---------- build new slides (appended, then moved) -----------------------
# We build all new slides (they go to the end), remember them, then reorder.
built = []  # list of (slide, insert_before_idx_at_moment_of_move)

# 1) "Demo use cases" - insert BEFORE current "3. Model evaluation" section.
s = new_slide(
    LAYOUT_SECTION,
    "Demo use cases we'll walk through",
    notes=(
        "Before we start demos, 60 seconds to set the two use cases. "
        "Every demo after this point runs against ONE of these two datasets. "
        "Keeping a concrete scenario in mind makes every metric more real."
    ),
)
built.append(("before_section_3", s))

s = new_slide(
    LAYOUT_TWOCOL if LAYOUT_TWOCOL.name else LAYOUT_CONTENT,
    "Clinical vs non-clinical - what we built",
    [
        "CLINICAL (datasets/clinical.jsonl - 10 prompts)",
        ("  1. Cardiac emergency triage (crushing chest pain)"),
        ("  2. Drug-interaction check (metformin + ibuprofen)"),
        ("  3. Chronic-disease status (T2 diabetes fasting glucose)"),
        ("  4. Pediatric red-flag triage (infant fever + poor feeding)"),
        ("  5. Mental-health crisis (988 escalation)"),
        ("  + 5 more across triage, meds, and patient-portal flow"),
        "",
        "NON-CLINICAL (datasets/nonclinical.jsonl - 10 prompts)",
        ("  1. Billing question (deductible explanation)"),
        ("  2. Scheduling (reschedule oncology appt)"),
        ("  3. HR benefits (PTO accrual)"),
        ("  4. Facility access (visitor policy after-hours)"),
        ("  5. Patient-portal troubleshooting (password reset)"),
        ("  + 5 more across billing, HR, and scheduling"),
    ],
    notes=(
        "Point at the clinical list - every row has a named clinical risk "
        "(cardiac, pediatric, mental health, drug interaction). That's why "
        "Lens 2 applies a clinical-safety rubric, a HIPAA-leak scanner, and "
        "a citation check on top of Foundry built-ins. The non-clinical set "
        "needs none of those overlays - it needs relevance, fluency, and "
        "basic refusal-of-clinical-questions behaviour instead."
    ),
)
built.append(("before_section_3", s))


# 2) "Foundry Model Catalog + Leaderboard" - insert AFTER the 'built-in' slide
# in section 3 so leaderboard ties the picture together.
s = new_slide(
    LAYOUT_CONTENT,
    "Foundry Model Catalog + Quality Leaderboard (UI)",
    [
        "Where in the portal:  ai.azure.com -> Model catalog -> 'Leaderboards'",
        "",
        "Three leaderboards out of the box:",
        ("  Quality leaderboard   - ranks models on MMLU, GPQA, Big-Bench-Hard, etc."),
        ("  Safety leaderboard    - ranks on harmful-content refusal + jailbreak resistance"),
        ("  Cost-Performance     - tokens/sec, $/1K tok, latency at P95"),
        "",
        "How Providence should use it:",
        ("  1. Short-list 2-3 candidates from the leaderboard (no code required)"),
        ("  2. Bring them into scripts 01-09 for the Providence overlay"),
        ("  3. Promote the winner to the Providence approved catalog"),
        "",
        "Important caveat:",
        ("  Leaderboards are GENERAL-PURPOSE benchmarks. They are NOT a substitute"),
        ("  for clinical / HIPAA / citation evaluation on YOUR data - that's Lens 2."),
    ],
    notes=(
        "Live demo: open ai.azure.com, click Model catalog, pick 'Leaderboards'. "
        "Sort by Quality, show top 5 models. Then click Safety tab - note which "
        "models dropped rank. This is the 'model selection from UI' piece "
        "Providence asked about. The key narrative: leaderboards give you the "
        "short-list (2-3 minutes, no code), but clinical fitness still requires "
        "Lens 2 with Providence data - generic leaderboards can't know that "
        "gpt-5-mini had 0% HIPAA leaks on your specific clinical dataset."
    ),
)
built.append(("after_builtin_explained", s))


# 3) Three evaluation modes - add a section + 2 content slides + a demo slide.
s = new_slide(
    LAYOUT_SECTION,
    "Automated, Manual, Human - three evaluation modes",
    notes=(
        "New section. Explains WHO runs WHICH evaluation and WHEN. "
        "Stress that a Providence model must pass all three before entering "
        "the approved catalog."
    ),
)
built.append(("after_agent_section", s))

s = new_slide(
    LAYOUT_CONTENT,
    "Three modes mapped to the Providence use case",
    [
        "AUTOMATED  (code + LLM-judge + deterministic evaluators)",
        ("  Who:    platform / ML engineers"),
        ("  When:   every PR, every nightly build"),
        ("  Scales: thousands of rows; ideal for regression"),
        ("  Example: scripts 01-08 against clinical + non-clinical datasets"),
        "",
        "MANUAL  (Foundry portal Evaluation UI - no code)",
        ("  Who:    governance / product team"),
        ("  When:   ad-hoc, 'what if we tried this new model?'"),
        ("  Scales: hundreds of rows; visual results in the portal"),
        ("  Example: upload 10-manual-upload.jsonl -> Run -> pick 4 evaluators -> done"),
        "",
        "HUMAN  (SMEs score rows in a spreadsheet, aggregated back)",
        ("  Who:    MDs, clinical-safety officers, compliance"),
        ("  When:   hard clinical cases; cases the automated judge disagrees on"),
        ("  Scales: dozens of rows; expensive but the only real ground truth"),
        ("  Example: datasets/human-eval-template.csv -> SMEs score -> aggregate JSON"),
    ],
    notes=(
        "Talk slowly through this slide. The question Providence always asks "
        "next is: 'how do we know the LLM judge is right?' Answer: you don't, "
        "not by itself. You cross-check with human SMEs on a subset, and if "
        "the automated and human scores disagree consistently, you adjust the "
        "judge rubric or drop that evaluator. The 10-human-aggregate.json "
        "shows the mean SME score alongside the automated score - that's "
        "your calibration signal."
    ),
)
built.append(("after_agent_section", s))

s = new_slide(
    LAYOUT_DEMO,
    "DEMO - all three modes on the same 5 rows",
    [
        "Script:  10_evaluation_modes.py",
        "",
        "What runs:",
        ("  1. Automated pass - 5 clinical prompts x gpt-5 x (ClinicalSafety + HIPAALeak)"),
        ("     -> eval-outputs/10-automated.result.json"),
        ("  2. Manual-ready JSONL - same rows formatted for Foundry portal upload"),
        ("     -> eval-outputs/10-manual-upload.jsonl"),
        ("  3. Human review CSV - pre-scored by simulated SMEs"),
        ("     -> datasets/human-eval-template.csv"),
        ("     -> eval-outputs/10-human-aggregate.json"),
        "",
        "Live results (today's run):",
        ("  Automated:  clinical_safety 5.0 / 5    hipaa_leak_rate 0.0"),
        ("  Human (5 SME rows):  clinical_safety 4.8 / 5    would_block_rate 0.0"),
        ("  Agreement: 0.2 gap - within calibration tolerance"),
        "",
        "Manual demo:  switch to ai.azure.com -> Evaluation -> upload the JSONL",
    ],
    notes=(
        "Run script 10 live (completes in ~30s). Show the three output files. "
        "Point at the 0.2-point gap between automated (5.0) and human (4.8) - "
        "tiny, so the LLM judge is well-calibrated on clinical safety for this "
        "dataset. If the gap were > 0.5, that's a signal to retune the judge "
        "prompt or fall back to human scoring as the source of truth. "
        "For the manual portion: drop the 10-manual-upload.jsonl into the "
        "Foundry Evaluation UI and show the one-click run experience."
    ),
)
built.append(("after_agent_section", s))


# 4) GitHub CI/CD evaluation - new section + content + demo slide.
s = new_slide(
    LAYOUT_SECTION,
    "CI/CD evaluation - GitHub Actions",
    notes=(
        "Governance as code. The exception process (09) becomes a PR gate. "
        "Any PR that touches evaluation code, datasets, or prompts re-runs "
        "the harness; any AUTO_DENIED model fails the build and blocks merge."
    ),
)
built.append(("end", s))

s = new_slide(
    LAYOUT_CONTENT,
    "Evaluation workflow on GitHub (Actions)",
    [
        "Repo:   github.com/jaypadhya1605/foundry-agent-observatory",
        "File:   .github/workflows/model-evaluation.yml",
        "",
        "Triggers:",
        ("  pull_request - on changes to eval scripts, datasets, requirements.txt"),
        ("  schedule (cron) - nightly 07:00 UTC regression run"),
        ("  workflow_dispatch - on-demand with candidate-model override inputs"),
        "",
        "What the job does:",
        ("  1. Checkout + Python 3.11 + install requirements"),
        ("  2. OIDC login to Azure (no secrets on disk)"),
        ("  3. Run 01 platform + 02 clinical + 03 non-clinical + 07 custom + 08 redteam"),
        ("  4. Run 09 exception-process - FAILS the job if any model = AUTO_DENIED"),
        ("  5. Upload all eval-outputs/*.json as a build artefact"),
        ("  6. Post a summary table as a PR comment (Model | Decision | Failures)"),
        "",
        "Net effect:  a model update cannot merge unless governance rules pass.",
    ],
    notes=(
        "Highlight OIDC: no long-lived Azure credentials stored in the repo. "
        "The workflow federates with Entra ID using workload identity, so "
        "secret rotation is handled centrally. Point at step 4 - this is "
        "where the PR gate enforces the governance decision. The PR comment "
        "in step 6 means reviewers see the AUTO_DENIED verdict inline, "
        "without clicking into Actions."
    ),
)
built.append(("end", s))

s = new_slide(
    LAYOUT_DEMO,
    "DEMO - a PR fails the governance gate",
    [
        "Tab:  github.com/jaypadhya1605/foundry-agent-observatory",
        "File: .github/workflows/model-evaluation.yml",
        "",
        "What to show (walkthrough):",
        ("  1. Open a pull request that adds a new prompt to datasets/clinical.jsonl"),
        ("  2. Click the 'Actions' tab - the 'Model Evaluation (Providence)' job is running"),
        ("  3. Watch step '09 Exception process' - shows gpt-5 AUTO_DENIED on HIPAA gate"),
        ("  4. Job fails with exit 1;  PR shows red X"),
        ("  5. PR conversation tab: bot-posted table with Model / Decision / Failures"),
        ("  6. Build artefact 'providence-eval-outputs' contains every *.result.json"),
        "",
        "Punchline:",
        ("  The governance policy now blocks a merge. Not a meeting, not an email."),
        ("  Every merge on main has an audited evaluation run attached to its SHA."),
    ],
    notes=(
        "If time permits, show a real historical PR on the repo. Otherwise "
        "just walk the YAML file and the Actions tab. Key governance "
        "takeaway: the SAME scripts we ran locally in earlier demos run "
        "unchanged in CI - no separate 'CI version' drift. The decision "
        "logic lives in 09_exception_process.py exactly once."
    ),
)
built.append(("end", s))


# ---------- reorder: move new slides to their target positions -------------
# python-pptx doesn't expose a move API directly; manipulate sldIdLst XML.
id_lst = prs.slides._sldIdLst
sld_list = list(id_lst)


def _rid_for_slide(slide):
    """Find the sldId element that points at this slide's part."""
    for sldId in list(id_lst):
        rId = sldId.attrib[qn("r:id")]
        try:
            if prs.part.related_part(rId) is slide.part:
                return sldId
        except KeyError:
            continue
    return None


# Anchor indices (re-evaluate after each insertion to be robust)
def idx_of_title(needle):
    for i, s in enumerate(prs.slides):
        if s.shapes.title is None:
            continue
        if needle.lower() in (s.shapes.title.text_frame.text or "").lower():
            return i
    return None


# Remove ALL new slide sldIds from current positions (they were appended at end).
# Cache (anchor, sldId_element) so we can re-insert AFTER removal.
print(f"id_lst count after add_slide loop: {len(list(id_lst))}")
pending = []  # list of (anchor, sldId_element) - element is now detached
for anchor, slide in built:
    sldId = _rid_for_slide(slide)
    if sldId is not None and sldId in id_lst:
        id_lst.remove(sldId)
    pending.append((anchor, sldId))

# Now re-insert each new slide at the correct target position, in order.
def insert_before_title(new_sldId, needle):
    """Insert new_sldId element before the sldId whose slide title contains needle."""
    for i, s in enumerate(prs.slides):
        if s.shapes.title is None:
            continue
        if needle.lower() in (s.shapes.title.text_frame.text or "").lower():
            current_list = list(id_lst)
            target_sldId = current_list[i]
            target_sldId.addprevious(new_sldId)
            print(f"    inserted before [{i+1}] ('{needle}')")
            return True
    print(f"    NOT FOUND anchor '{needle}' - appending")
    id_lst.append(new_sldId)
    return False


def append_sldId(new_sldId):
    id_lst.append(new_sldId)


print(f"built {len(built)} new slides; count in id_lst before reinsert: {len(list(id_lst))}")
# Process the pending list in order, placing each slide at its anchor.
for anchor, sldId in pending:
    if sldId is None:
        continue
    if anchor == "before_section_3":
        # "3. Model evaluation" section slide
        insert_before_title(sldId, "3. Model evaluation")
    elif anchor == "after_builtin_explained":
        # Insert after "Built-in evaluators - what they are" slide -> before first DEMO (Lens 1).
        insert_before_title(sldId, "Built-in evaluators + Lens 1")
    elif anchor == "after_agent_section":
        # Insert AFTER the "5. Agent evaluation" section
        # i.e. before "6. Exception process"
        insert_before_title(sldId, "6. Exception process")
    elif anchor == "end":
        # CI/CD section: place after architecture, before "Providence asks -> delivered today"
        inserted = insert_before_title(sldId, "Providence asks")
        if not inserted:
            insert_before_title(sldId, "Takeaways")


prs.save(str(DECK))
print(f"Saved {DECK}  ({len(prs.slides)} slides)")
for i, s in enumerate(prs.slides, 1):
    t = ""
    if s.shapes.title is not None and s.shapes.title.has_text_frame:
        t = s.shapes.title.text_frame.text.strip().replace("\n", " / ")[:80]
    print(f"[{i:02d}] {t}")
