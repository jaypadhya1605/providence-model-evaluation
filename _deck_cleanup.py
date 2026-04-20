"""Final cleanup: delete duplicate slide 28, scrub script names from slide 37."""
import sys
from pathlib import Path
sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

DECK = Path("Providence - Model Evaluation.pptx")
prs = Presentation(str(DECK))

# ---- delete the duplicate "Three evaluation modes" slide (the former section)
# There are two slides with the same title; keep the second (content).
id_lst = prs.slides._sldIdLst
current = list(id_lst)
titles = []
for sldId in current:
    rId = sldId.attrib[qn("r:id")]
    slide = prs.part.related_part(rId)
    titles.append(slide)

# Iterate slides and find duplicate title pair
seen = {}
dup_to_remove = None
for i, s in enumerate(prs.slides):
    if s.shapes.title and s.shapes.title.has_text_frame:
        t = s.shapes.title.text_frame.text.strip()
        if "Three evaluation modes" in t:
            if "three_modes" in seen:
                # remove the FIRST one (section placeholder), keep the content one
                dup_to_remove = seen["three_modes"]
                break
            seen["three_modes"] = i

if dup_to_remove is not None:
    sldId_elt = current[dup_to_remove]
    # also drop the relationship + the part from the presentation
    rId = sldId_elt.attrib[qn("r:id")]
    try:
        prs.part.drop_rel(rId)
    except Exception:
        pass
    id_lst.remove(sldId_elt)
    print(f"Removed duplicate slide at index {dup_to_remove + 1}")

# ---- scrub script names from slide 'Providence asks -> delivered today'
C_BLUE = RGBColor(0x12, 0x5E, 0xA6)
C_TEXT = RGBColor(0x25, 0x32, 0x3A)

REPLACEMENTS = [
    ("✔  Clinical vs non-clinical evaluation       (scripts 02 + 03)",
     "✔  Clinical vs non-clinical evaluation"),
    ("✔  Platform approval vs use-case separation   (script 01 vs 02/03)",
     "✔  Platform approval vs use-case separation"),
    ("✔  Built-in evaluators explained + efficacy   (script 04)",
     "✔  Built-in evaluators explained + efficacy"),
    ("✔  Exception process demo-ready               (script 09)",
     "✔  Exception process demo-ready"),
    ("✔  Providence-specific guardrail overlay      (_custom_evaluators.py + script 07)",
     "✔  Providence-specific guardrail overlay"),
    ("✔  Agentic evaluation on NEW Foundry agents   (scripts 05 + 06)",
     "✔  Agentic evaluation on NEW Foundry agents"),
    ("✔  Red-team / adversarial harness             (script 08)",
     "✔  Red-team / adversarial harness"),
]

for s in prs.slides:
    if s.shapes.title and "Providence asks" in (s.shapes.title.text_frame.text or ""):
        for shp in s.shapes:
            if not shp.has_text_frame or shp is s.shapes.title:
                continue
            for para in shp.text_frame.paragraphs:
                for run in para.runs:
                    txt = run.text
                    for old, new in REPLACEMENTS:
                        if old.strip() in txt:
                            txt = new
                    # regex-style scrub: anything with "(script" in it
                    if "(script" in txt or "_custom_evaluators.py" in txt:
                        # take everything up to the first "(" and strip
                        import re
                        txt = re.sub(r"\s*\([^)]*(?:script|\.py)[^)]*\)\s*", "", txt).rstrip()
                    run.text = txt
        # add extra detail to speaker notes instead
        notes_tf = s.notes_slide.notes_text_frame
        existing = notes_tf.text
        extra = (
            "\n\nScript map (for reference, not for the slide):\n"
            "  01_platform_approval_eval.py   - Lens 1\n"
            "  02_usecase_clinical_eval.py    - Lens 2 clinical\n"
            "  03_usecase_nonclinical_eval.py - Lens 2 non-clinical\n"
            "  04_builtin_evaluators_explained.py - what each built-in returns\n"
            "  05_create_foundry_agent.py + 06_agent_evaluation.py - agents\n"
            "  07_custom_evaluators.py + _custom_evaluators.py - Providence overlay\n"
            "  08_redteam_adversarial.py - red-team\n"
            "  09_exception_process.py - final governance decision\n"
            "  10_evaluation_modes.py - automated / manual / human\n"
            "  .github/workflows/model-evaluation.yml - CI gate"
        )
        if "Script map" not in existing:
            notes_tf.text = existing + extra
        break

prs.save(str(DECK))
print(f"Saved {DECK} ({len(prs.slides)} slides)")
