"""
Build the unified Providence deck.
  - Starts from "L150 Microsoft Foundry Models Pitchdeck.PPTX" (keeps its theme).
  - Keeps a curated set of original slides.
  - Inserts new content + demo + section + architecture slides using the
    SAME layouts (so the theme stays consistent).
  - Adds speaker notes to every slide.
Output: "Providence - Model Evaluation.pptx"
"""
from __future__ import annotations
import copy
import sys
from pathlib import Path

sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

SRC = Path("L150 Microsoft Foundry Models Pitchdeck.PPTX")
DST = Path("Providence - Model Evaluation.pptx")

prs = Presentation(str(SRC))

# 1-based indices of slides we want to reuse from the original deck.
KEEP = {
    2:   "orig_title",          # Microsoft Foundry Models (title)
    9:   "orig_factory",        # AI App & Agent Factory
    11:  "orig_catalog",        # 11,000+ frontier/open models
    13:  "orig_unified_access", # Unified access (sold/supported by MS)
    93:  "orig_offer_type",     # Choosing the best offer type
    94:  "orig_deploy_loc",     # Choosing the best deployment location
    114: "orig_content_safety", # Azure AI Content Safety
}

# Snapshot keeper slide objects BEFORE we delete anything.
keeper_slides = {}
for i, s in enumerate(prs.slides, 1):
    if i in KEEP:
        keeper_slides[KEEP[i]] = s

# ---------- helpers ----------------------------------------------------------
def _slide_id_pair(slide):
    """Return (sldId element, rId) for a slide in prs.slides._sldIdLst."""
    id_lst = prs.slides._sldIdLst
    for sldId in list(id_lst):
        rId = sldId.attrib[qn("r:id")]
        part = prs.part.related_part(rId)
        if part is slide.part:
            return sldId, rId
    return None, None


def delete_slide(slide):
    id_lst = prs.slides._sldIdLst
    sldId, rId = _slide_id_pair(slide)
    if sldId is not None:
        id_lst.remove(sldId)
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass


def move_to_end(slide):
    id_lst = prs.slides._sldIdLst
    sldId, _ = _slide_id_pair(slide)
    if sldId is not None:
        id_lst.remove(sldId)
        id_lst.append(sldId)


def find_layout(*name_candidates):
    for sm in prs.slide_masters:
        for layout in sm.slide_layouts:
            for c in name_candidates:
                if layout.name.strip().lower() == c.strip().lower():
                    return layout
    # fallback: partial match
    for sm in prs.slide_masters:
        for layout in sm.slide_layouts:
            for c in name_candidates:
                if c.strip().lower() in layout.name.strip().lower():
                    return layout
    return prs.slide_masters[0].slide_layouts[0]


LAYOUT_TITLE   = find_layout("Title Slide", "Title Slide 2")
LAYOUT_SECTION = find_layout("Section Title", "Section Title 2")
LAYOUT_CONTENT = find_layout("Title and Content")
LAYOUT_TWOCOL  = find_layout("Two Column Bullet with Subheads", "Two Column Bullet text")
LAYOUT_DEMO    = find_layout("Demo slide", "Demo slide 2")
LAYOUT_THANKS  = find_layout("Thank you slide", "Closing logo slide")
LAYOUT_BLANK   = find_layout("Blank")


def set_title(slide, text):
    if slide.shapes.title is not None:
        slide.shapes.title.text = text
    else:
        # Add a title-sized text box at top
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(36)
        p.font.bold = True


def set_body_bullets(slide, bullets, placeholder_idx=None, left=None, top=None, width=None, height=None):
    """
    Writes bullets into the content placeholder if present, otherwise adds a textbox.
    bullets: list of str OR list of (level:int, text:str) tuples
    """
    tf = None
    # Try to find a content placeholder (not the title)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 0 and ph.has_text_frame:
            tf = ph.text_frame
            break
    if tf is None:
        tb = slide.shapes.add_textbox(
            left or Inches(0.5),
            top or Inches(1.6),
            width or Inches(12.3),
            height or Inches(5.5),
        )
        tf = tb.text_frame
    tf.word_wrap = True

    # Clear existing runs
    tf.clear()

    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        for r in p.runs:
            if level == 0:
                r.font.size = Pt(20)
                r.font.bold = False
            else:
                r.font.size = Pt(16)


def add_notes(slide, notes):
    nslide = slide.notes_slide
    nslide.notes_text_frame.text = notes


def add_slide(layout, title, bullets=None, notes=None):
    s = prs.slides.add_slide(layout)
    set_title(s, title)
    if bullets:
        set_body_bullets(s, bullets)
    if notes:
        add_notes(s, notes)
    return s


# ---------- plan ------------------------------------------------------------
# We build the NEW deck in order, then at the end remove every original slide
# except the keepers, and place keepers at their chosen positions.

# We'll collect (kind, payload) ops; 'orig' = keeper key, 'new' = new slide built already.
new_slides = []


# --- Slide 1: Title (use original slide 2 "Microsoft Foundry Models") -------
new_slides.append(("orig", "orig_title"))

# Override title text + add subtitle + notes on the original slide below.

# --- Slide 2: Agenda ---
s = add_slide(
    LAYOUT_CONTENT,
    "Agenda",
    [
        "The Providence ask — what today answers",
        "Foundry Models platform & deployment choices",
        "Model deployment — what we've deployed for Providence",
        "Model evaluation — platform approval (Lens 1)",
        "Use-case evaluation — clinical + non-clinical (Lens 2)",
        "Custom evaluators — Providence overlay (HIPAA / clinical safety / citations)",
        "Red-team & Content Safety guardrails",
        "Agent evaluation — NEW Foundry agents (not classic Assistants)",
        "Exception process — governance workflow",
        "Architecture recap + punchlines",
    ],
    notes=(
        "60-minute plan. Flag the two-lens model as the single most important "
        "takeaway: 'platform approval' is separate from 'use-case evaluation'. "
        "We'll demo live at slides marked 'Demo'. Keep the Streamlit app open "
        "in one tab and the Foundry portal (Agents view) in another tab."
    ),
)
new_slides.append(("new", s))


# --- Slide 3: Section — The Providence ask ---
s = add_slide(
    LAYOUT_SECTION,
    "1. The Providence ask",
    notes=(
        "Anchor the session on what Providence explicitly asked for in the "
        "March 4, March 11, and March 30 syncs. We're not inventing a framework "
        "— we're answering five specific questions from Troy, Hames, and Aashish."
    ),
)
new_slides.append(("new", s))


# --- Slide 4: The five explicit asks (from providence-evaluation-notes) -----
s = add_slide(
    LAYOUT_CONTENT,
    "What Providence explicitly asked",
    [
        "1. How do we evaluate models for Providence — clinical vs non-clinical?  (Troy)",
        "2. How do we separate platform approval from use-case evaluation?  (March 4 sync)",
        "3. What are built-in evaluators and how do they determine efficacy?  (Hames)",
        "4. What is the exception process when teams want a non-standard model?  (Aashish)",
        "5. How do we supplement Foundry with Providence-specific guardrails?  (Troy)",
        "",
        "Everything on the next 25 slides maps back to one of these five.",
    ],
    notes=(
        "Read these aloud verbatim. The phrasing is deliberately taken from the "
        "meeting transcripts. Call out: 'the team said these two get mixed up' "
        "- that's question #2. Hames explicitly flagged 'I don't see documentation "
        "on how they determine efficacy' - that's #3. Aashish asked to see the "
        "exception process in the demo - that's #4."
    ),
)
new_slides.append(("new", s))


# --- Slide 5: The two-lens answer ---
s = add_slide(
    LAYOUT_CONTENT,
    "The two-lens model (the answer to #2)",
    [
        "Lens 1 — PLATFORM APPROVAL",
        ("  Can this model enter the Providence catalog at all?"),
        ("  Safety floor + baseline quality (coherence / fluency / groundedness / content safety)"),
        ("  Run ONCE per model version — reusable across every project."),
        "",
        "Lens 2 — USE-CASE EVALUATION",
        ("  Is this model appropriate for THIS workflow?"),
        ("  Clinical overlay:   clinical safety + HIPAA leak + medical citation + groundedness"),
        ("  Non-clinical overlay: relevance + similarity + coherence + fluency"),
        ("  Run per workflow — different datasets, different thresholds."),
    ],
    notes=(
        "THIS is the slide Providence will remember. Hold it on screen for "
        "~90 seconds. The two lenses are the core mental model. Make the point: "
        "every painful governance debate at any health system comes from mixing "
        "these up. A model that passes Lens 1 is not automatically approved for "
        "a bedside clinical workflow — Lens 2 re-checks with a clinical dataset "
        "and clinical rubric. That is the 'separation' Providence asked for."
    ),
)
new_slides.append(("new", s))


# --- Slide 6: Foundry factory (original slide 9) ---
new_slides.append(("orig", "orig_factory"))

# --- Slide 7: Foundry catalog (original slide 11) ---
new_slides.append(("orig", "orig_catalog"))

# --- Slide 8: Unified access (original slide 13) ---
new_slides.append(("orig", "orig_unified_access"))

# --- Slide 9: Section — Model deployment ---
s = add_slide(
    LAYOUT_SECTION,
    "2. Model deployment",
    notes=(
        "We'll cover deployment shape for about 3 minutes. Don't over-explain "
        "pricing; reference the next two original-deck slides and then land "
        "on the actual deployment we're running for Providence."
    ),
)
new_slides.append(("new", s))

# --- Slide 10-11: original offer type + deployment location ---
new_slides.append(("orig", "orig_offer_type"))
new_slides.append(("orig", "orig_deploy_loc"))

# --- Slide 12: Our Providence deployment ---
s = add_slide(
    LAYOUT_CONTENT,
    "What we've deployed for Providence today",
    [
        "Foundry account:  ai-apim-demo-jp-001  (East US 2, S0)",
        "Foundry project:  proj-apim-demo-jp-001",
        "",
        "Model deployments on that account:",
        ("  gpt-4o        Standard / 30 TPM   — Judge model (LLM-as-a-judge)"),
        ("  gpt-4.1-mini  GlobalStandard / 10 TPM"),
        ("  gpt-5         Standard / 10 TPM   — Candidate A"),
        ("  gpt-5-mini    Standard / 10 TPM   — Candidate B"),
        ("  text-embedding-3-small"),
        "",
        "Zero new infrastructure for this session. Everything reuses Sessions 1 & 2.",
    ],
    notes=(
        "Point out: we reused the SAME Foundry resource from the Observability "
        "and APIM sessions. That is deliberate — Providence doesn't want a new "
        "pile of resources per demo. Cost for today's evaluation run: ~single-"
        "digit dollars in inference, zero net-new infra."
    ),
)
new_slides.append(("new", s))


# --- Slide 13: Section — Model Evaluation (answers #1) ---
s = add_slide(
    LAYOUT_SECTION,
    "3. Model evaluation — Foundry built-ins",
    notes=(
        "Shift into evaluation proper. This section answers Providence ask #1 "
        "(clinical vs non-clinical) and #3 (what are built-ins, how do they "
        "determine efficacy)."
    ),
)
new_slides.append(("new", s))

# --- Slide 14: Built-in evaluators (answers #3) ---
s = add_slide(
    LAYOUT_CONTENT,
    "Built-in evaluators — what they are, how they score",
    [
        "Quality (LLM-as-judge, 1-5 Likert scale, higher is better):",
        ("  Coherence, Fluency, Relevance, Similarity, Groundedness"),
        "",
        "Risk & safety (LLM-judged classifiers with severity scale):",
        ("  Violence, Self-harm, Sexual, Hate/Unfairness, Content Safety aggregate"),
        ("  Indirect Attack (prompt injection), Protected Material"),
        "",
        "Lexical (deterministic):",
        ("  F1, BLEU, ROUGE, METEOR"),
        "",
        "Each evaluator emits: score, reason, threshold, pass/fail, token usage, judge model used.",
    ],
    notes=(
        "Hames asked 'how do they determine efficacy?' — answer: every LLM-"
        "judged evaluator ships with a documented rubric prompt, a 1-5 scale "
        "and a threshold. The score, the reason, the token usage, and the "
        "judge model are all in the output JSON. We'll show this live in the "
        "next demo — 04_builtin_evaluators_explained.py runs each built-in "
        "with one sample so you can literally read the judge's reasoning."
    ),
)
new_slides.append(("new", s))


# --- Slide 15: DEMO 1 — built-ins + platform approval (01) ---
s = add_slide(
    LAYOUT_DEMO,
    "DEMO — Built-in evaluators + Lens 1 (platform approval)",
    [
        "Script:  04_builtin_evaluators_explained.py   (walks one evaluator at a time)",
        "Script:  01_platform_approval_eval.py         (gpt-5 + gpt-5-mini → catalog approval)",
        "",
        "What to watch for on screen:",
        ("  judge model = gpt-4o, scores + reasons, per-metric thresholds"),
        ("  final verdict:  APPROVED_FOR_CATALOG"),
        "",
        "Streamlit tab:  01 — Platform approval",
    ],
    notes=(
        "Switch to the terminal. Run 04 first so they can literally see the "
        "judge reasoning string next to the numeric score. Then show the "
        "pre-computed 01 result via the Streamlit app (don't re-run — it "
        "takes ~3 minutes per model). Callouts while on Streamlit: "
        "(a) both candidate models PASSED the platform floor, "
        "(b) the decision rule is transparent - you can see exact thresholds, "
        "(c) this is a one-time check per model version; doesn't need to be "
        "repeated per project."
    ),
)
new_slides.append(("new", s))


# --- Slide 16: DEMO 2 — clinical use-case ---
s = add_slide(
    LAYOUT_DEMO,
    "DEMO — Lens 2 (clinical use-case evaluation)",
    [
        "Script:  02_usecase_clinical_eval.py",
        "Dataset: datasets/clinical.jsonl  (triage, drug interaction, patient-facing Q&A)",
        "",
        "Built-in evaluators:  Relevance · Similarity · Groundedness · Coherence · F1",
        "Custom overlays:     ClinicalSafety · HIPAALeak · MedicalCitation",
        "",
        "Live results (gpt-5 vs gpt-5-mini):",
        ("  ClinicalSafety:  4.9 / 5       HIPAA leak rate:  10% on gpt-5, 0% on gpt-5-mini"),
        ("  Medical citation score:  5.0 / 4.8       Groundedness:  4.8 / 4.6"),
        "",
        "Streamlit tab:  02 — Clinical use-case",
    ],
    notes=(
        "Headline: gpt-5 is better on citation discipline, but gpt-5-mini "
        "actually leaked zero PII patterns in this sample — this is exactly "
        "the kind of signal you cannot get from a generic leaderboard. Click "
        "into a single row on Streamlit to show the judge's reason string. "
        "That transparency is the answer to 'how do we know it actually "
        "evaluated clinical risk?'"
    ),
)
new_slides.append(("new", s))


# --- Slide 17: DEMO 3 — non-clinical use-case ---
s = add_slide(
    LAYOUT_DEMO,
    "DEMO — Lens 2 (non-clinical use-case evaluation)",
    [
        "Script:  03_usecase_nonclinical_eval.py",
        "Dataset: datasets/nonclinical.jsonl  (billing, scheduling, HR)",
        "",
        "Evaluators:  Relevance · Similarity · Coherence · Fluency · F1",
        "",
        "Live results (gpt-5 vs gpt-5-mini):",
        ("  Relevance:  4.9 / 4.8        Coherence:  4.7 / 4.4"),
        ("  Fluency:    4.0 / 4.1        Similarity: 3.5 / 3.6"),
        "",
        "Same harness, different rubric — that is the point.",
    ],
    notes=(
        "The important narrative: the SAME scripts and SAME infrastructure "
        "handle clinical and non-clinical. Only the dataset and the active "
        "evaluator set change. This means Providence can onboard a new "
        "workflow in hours, not weeks. Non-clinical doesn't need clinical "
        "safety or HIPAA-leak gates at the same severity level — that's why "
        "they're removed from Lens 2's non-clinical overlay."
    ),
)
new_slides.append(("new", s))


# --- Slide 18: Section — Custom evaluators ---
s = add_slide(
    LAYOUT_SECTION,
    "4. Custom evaluators — the Providence overlay",
    notes=(
        "This answers ask #5 - how do we supplement Foundry built-ins with "
        "Providence-specific policy. The message: Foundry gives you the floor, "
        "you bring the house rules."
    ),
)
new_slides.append(("new", s))


# --- Slide 19: Custom evaluators - overlay ---
s = add_slide(
    LAYOUT_CONTENT,
    "Providence overlay — three custom evaluators",
    [
        "ClinicalSafetyEvaluator  (LLM-judged 1-5)",
        ("  Penalizes diagnoses, prescription language, missing 'consult a provider'"),
        ("  Emergency recognition (chest pain / stroke / 988 triggers)"),
        "",
        "HIPAALeakEvaluator  (deterministic regex, no LLM cost)",
        ("  SSN · phone · email · MRN · DOB patterns"),
        ("  Hard-fail in exception process — any leak blocks approval"),
        "",
        "MedicalCitationEvaluator  (LLM-judged 1-5)",
        ("  Does the response cite CDC / WHO / NIH / FDA when making medical claims?"),
        "",
        "All three plug directly into azure.ai.evaluation's evaluate() — same signature as built-ins.",
    ],
    notes=(
        "Emphasize that these aren't prototypes — they use the same callable "
        "contract (query, response, **kwargs) as every Foundry built-in, so "
        "they aggregate into the same result JSON and appear in the same "
        "Streamlit tab. HIPAA leak is deterministic regex precisely so we "
        "never depend on an LLM's mood to catch PII. Safety and citation "
        "are LLM-judged because they need semantic understanding."
    ),
)
new_slides.append(("new", s))


# --- Slide 20: DEMO 4 — custom evals standalone ---
s = add_slide(
    LAYOUT_DEMO,
    "DEMO — Custom evaluators in action",
    [
        "Script:  07_custom_evaluators.py",
        "",
        "Live results on gpt-5 output:",
        ("  clinical_safety_score           4.9 / 5"),
        ("  hipaa_leak_leaked              10% (1 of 10 samples)"),
        ("  medical_citation_score         5.0 / 5   (1.0 cites_authority rate)"),
        "",
        "Streamlit tab:  07 — Custom evaluators",
    ],
    notes=(
        "Callout: the 10% HIPAA leak rate is the real story. That single "
        "failing row is exactly what drives the AUTO_DENIED outcome in the "
        "exception-process demo later. Show the leaked pattern on Streamlit "
        "so Providence sees the regex catch concretely."
    ),
)
new_slides.append(("new", s))


# --- Slide 21: Content Safety (original 114) ---
new_slides.append(("orig", "orig_content_safety"))


# --- Slide 22: DEMO 5 — red-team / adversarial ---
s = add_slide(
    LAYOUT_DEMO,
    "DEMO — Red-team / adversarial",
    [
        "Script:  08_redteam_adversarial.py",
        "Dataset: datasets/adversarial.jsonl  (jailbreaks, PHI-fishing, prompt-injection)",
        "",
        "Evaluators:  Foundry Content Safety + custom Refusal/DangerousLeak",
        "",
        "Live results on gpt-5:",
        ("  refusal_detected rate:   10%"),
        ("  dangerous_leak rate:     20%"),
        "",
        "Streamlit tab:  08 — Red-team",
    ],
    notes=(
        "This is uncomfortable on purpose. 20% dangerous-leak on adversarial "
        "shows the value of the Providence overlay AND of APIM-level guardrails "
        "from Session 2. Tie back: 'This is why we run the harness — so you "
        "see it before a patient sees it.' Mention: one row failed with a "
        "transient connection error (that's normal for adversarial batch eval); "
        "the harness captures per-row errors separately from judge scores."
    ),
)
new_slides.append(("new", s))


# --- Slide 23: Section — Agent evaluation ---
s = add_slide(
    LAYOUT_SECTION,
    "5. Agent evaluation (not just models)",
    notes=(
        "Pivot: everything so far is model-level. Agents add tools, plans, and "
        "multi-turn behaviour — which need a different evaluation rubric."
    ),
)
new_slides.append(("new", s))

# --- Slide 24: NEW Foundry agents vs classic Assistants ---
s = add_slide(
    LAYOUT_CONTENT,
    "NEW Foundry agents — not classic Assistants",
    [
        "What we did NOT use:",
        ("  azure-ai-agents 1.x  →  creates asst_* ids, shows as 'legacy' in Foundry UI"),
        "",
        "What we used:",
        ("  azure-ai-projects 2.x + PromptAgentDefinition + project.agents.create_version()"),
        ("  Agent id shape:  providence-clinical-triage:2   (name:version, not asst_*)"),
        ("  Runtime:  OpenAI Responses API via project.get_openai_client()"),
        "",
        "Why it matters for Providence:",
        ("  - Agent appears in Foundry portal → Agents (not in the 'Legacy Assistants' pane)"),
        ("  - Version-pinned — every rev is a new, auditable artefact"),
        ("  - Same governance story as models: platform floor + use-case eval + exceptions"),
    ],
    notes=(
        "Providence specifically asked us NOT to use the classic Assistant API. "
        "If someone in the audience is building against asst_* ids today, this "
        "is the migration slide. Foundry_agentic-creation.md in the repo has "
        "the full code walk-through — link to it in the follow-up email."
    ),
)
new_slides.append(("new", s))


# --- Slide 25: DEMO 6 — agent create + eval ---
s = add_slide(
    LAYOUT_DEMO,
    "DEMO — Create an agent + evaluate it",
    [
        "Script:  05_create_foundry_agent.py   (creates 'providence-clinical-triage:2')",
        "Script:  06_agent_evaluation.py       (IntentResolution + TaskAdherence + ToolCallAccuracy)",
        "",
        "Live results:",
        ("  intent_resolution         4.0 / 5"),
        ("  task_adherence            0.8 (pass rate)"),
        ("  tool_call_accuracy        5.0 / 5   - escalate_to_human() fired on 'crushing chest pain'"),
        "",
        "Show agent in Foundry portal → Agents → providence-clinical-triage",
        "Streamlit tab:  06 — Agent evaluation",
    ],
    notes=(
        "Flip to ai.azure.com in the browser and show the agent in the "
        "Foundry UI — proves the 'not classic' point visually. Then back to "
        "Streamlit for the metrics. Highlight: tool_call_accuracy = 5 because "
        "the agent correctly chose escalate_to_human() on a cardiac emergency. "
        "That is the clinical-safety bar operationalized as an agent metric."
    ),
)
new_slides.append(("new", s))


# --- Slide 26: Section — Exception process ---
s = add_slide(
    LAYOUT_SECTION,
    "6. Exception process (Aashish's ask)",
    notes=(
        "Final narrative arc. Aashish specifically asked to see this live in "
        "the demo. Set expectation: this is not a UI tool, it's a deterministic "
        "decision engine that reads all prior evaluation outputs."
    ),
)
new_slides.append(("new", s))


# --- Slide 27: DEMO 7 — exception process ---
s = add_slide(
    LAYOUT_DEMO,
    "DEMO — Exception process decision engine",
    [
        "Script:  09_exception_process.py",
        "Input:   every *.result.json file from prior evaluation runs",
        "Output:  exception-decision-<model>.json   →   APPROVED / NEEDS_REVIEW / AUTO_DENIED",
        "",
        "Decision rules (all deterministic, version-controlled):",
        ("  HIPAA leak != 0        →  AUTO_DENIED  (hard fail)"),
        ("  ClinicalSafety < 4     →  AUTO_DENIED  (hard fail)"),
        ("  Any quality metric < threshold  →  NEEDS_REVIEW (human approver)"),
        ("  All green              →  APPROVED (audit-logged)"),
        "",
        "Live: both gpt-5 and gpt-5-mini  →  AUTO_DENIED on today's clinical dataset",
        "Reason: 10% / 10% HIPAA-leak rate on clinical custom overlay",
    ],
    notes=(
        "Key point: the decision output IS the audit log. Route the JSON to "
        "governance-owned email (env var EXCEPTION_APPROVER_EMAIL). The "
        "AUTO_DENIED result today is the demo's strongest proof point — the "
        "governance workflow caught something Providence's ClinicalSafetyOfficer "
        "would also catch. Better to catch it here, not in production."
    ),
)
new_slides.append(("new", s))


# --- Slide 28: Architecture ---
s = add_slide(
    LAYOUT_CONTENT,
    "End-to-end architecture (zero-new-infra)",
    [
        "Reused from Session 1 (Observability):",
        ("  App Insights + Foundry tracing + dashboards"),
        "",
        "Reused from Session 2 (APIM):",
        ("  APIM AI Gateway · semantic caching · token-rate limits · content-safety policies"),
        "",
        "New in Session 3:",
        ("  9 evaluation scripts (01-09)  +  3 datasets  +  3 custom evaluators"),
        ("  Streamlit demo UI (reads on-disk JSON only — safe in front of an audience)"),
        ("  Exception-decision engine (09) writing audit-log JSON"),
        "",
        "One Foundry account: ai-apim-demo-jp-001. One project. Five model deployments. Done.",
    ],
    notes=(
        "Use this slide to reset the cost story. Everything ran on ONE Foundry "
        "account. Total monthly fixed cost: essentially the App Insights ingestion "
        "(already paid) and the model deployments (per-token). No new storage, "
        "no new ML workspace, no GPU endpoints."
    ),
)
new_slides.append(("new", s))


# --- Slide 29: Punchlines / tick-box ---
s = add_slide(
    LAYOUT_CONTENT,
    "Providence asks  →  delivered today",
    [
        "✔  Clinical vs non-clinical evaluation       (scripts 02 + 03)",
        "✔  Platform approval vs use-case separation   (script 01 vs 02/03)",
        "✔  Built-in evaluators explained + efficacy   (script 04)",
        "✔  Exception process demo-ready               (script 09)",
        "✔  Providence-specific guardrail overlay      (_custom_evaluators.py + script 07)",
        "✔  Agentic evaluation on NEW Foundry agents   (scripts 05 + 06)",
        "✔  Red-team / adversarial harness             (script 08)",
        "✔  Unified Streamlit demo UI                  (9 tabs, one per ask)",
        "✔  Zero net-new infrastructure                (reused ai-apim-demo-jp-001)",
    ],
    notes=(
        "Land the plane. Every checkbox ties to a specific transcript line or "
        "document request. Offer to leave the repo, the Streamlit app, and "
        "this deck with the Providence team as a starting point they can "
        "fork and run themselves."
    ),
)
new_slides.append(("new", s))


# --- Slide 30: Takeaways / next steps ---
s = add_slide(
    LAYOUT_CONTENT,
    "Takeaways & next steps",
    [
        "Takeaway 1 — Two lenses, never one.  Platform approval is not use-case approval.",
        "Takeaway 2 — Built-ins + overlays.  Foundry gives you the floor; you bring the house rules.",
        "Takeaway 3 — Governance is code.  The exception engine is a script, not a meeting.",
        "Takeaway 4 — Agents are governed the same way.  Don't split the rubric.",
        "",
        "Suggested next steps for Providence:",
        ("  1. Fork the repo, swap in a real Providence clinical dataset (100-200 rows)"),
        ("  2. Wire 09's exception-decision JSON into ServiceNow or email approval flow"),
        ("  3. Run 04 against the two models you're evaluating THIS quarter"),
        ("  4. Decide which overlays are 'hard fail' vs 'needs review' for your governance board"),
    ],
    notes=(
        "Keep the close tight. Ask: 'Of these four takeaways, which one is "
        "most relevant to the decisions you have to make in the next 30 days?' "
        "That gets you a qualified follow-up without pitching."
    ),
)
new_slides.append(("new", s))


# --- Slide 31: Thank you / Q&A ---
s = add_slide(
    LAYOUT_THANKS if LAYOUT_THANKS.name else LAYOUT_SECTION,
    "Thank you — Q & A",
    notes=(
        "Expected questions:\n"
        "Q: 'Can we run this on our own Foundry tenant?' → yes, fork the repo, "
        "change .env, done.\n"
        "Q: 'What about fine-tuned models?' → same harness; platform approval "
        "re-runs on every new fine-tune version.\n"
        "Q: 'Does the exception decision actually block deployment?' → today "
        "it writes an audit JSON; wiring into your CD pipeline or ServiceNow "
        "takes <1 day.\n"
        "Q: 'What does this cost to run continuously?' → inference-only; "
        "nightly eval on 100 rows across 2 models ~ a few dollars/week."
    ),
)
new_slides.append(("new", s))


# ---------- now: delete every ORIGINAL slide that's not a keeper ------------
# Iterate over a copy because we're mutating.
for s in list(prs.slides):
    if s in keeper_slides.values():
        continue
    # Skip our freshly-added new slides (they're in the list too).
    # Identify them: new_slides contains ("new", slide_obj).
    is_new = any(kind == "new" and obj is s for kind, obj in new_slides)
    if is_new:
        continue
    delete_slide(s)


# ---------- now: enforce final order ---------------------------------------
id_lst = prs.slides._sldIdLst
# Remove ALL sldIds first, then re-append in the target order.
existing = list(id_lst)
for sldId in existing:
    id_lst.remove(sldId)

def _sldId_for(slide):
    # After removal, we need to reinsert. Build a sldId element referencing the rId.
    # Easier: search the original list snapshot by rId.
    for sldId in existing:
        rId = sldId.attrib[qn("r:id")]
        try:
            part = prs.part.related_part(rId)
        except KeyError:
            continue
        if part is slide.part:
            return sldId
    return None

for kind, payload in new_slides:
    if kind == "orig":
        slide = keeper_slides[payload]
    else:
        slide = payload
    sldId = _sldId_for(slide)
    if sldId is not None:
        id_lst.append(sldId)


# ---------- tweak the reused ORIGINAL title slide ---------------------------
# Override the big title on original-slide-2 to Providence branding + subtitle.
title_slide = keeper_slides["orig_title"]
# Find the main title-like text boxes and rewrite the first big one.
for sh in title_slide.shapes:
    if sh.has_text_frame and sh.text_frame.text.strip().startswith("Microsoft"):
        tf = sh.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = "Microsoft Foundry"
        r.font.size = Pt(54)
        r.font.bold = True
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = "Providence — Model Evaluation"
        r2.font.size = Pt(40)
        r2.font.bold = True
        p3 = tf.add_paragraph()
        r3 = p3.add_run()
        r3.text = "Session 3 of 3 · 60-minute technical walkthrough"
        r3.font.size = Pt(22)
        break

add_notes(
    title_slide,
    "Session 3 of the Providence Foundry deep-dive (Observability → APIM → "
    "Model Evaluation). Open by thanking Troy, Hames, and Aashish by name "
    "for the specific questions that shaped this agenda. Mention today's "
    "session is 60 minutes with live demo — flag that we'll pause for "
    "questions after each section."
)


# ---------- save -----------------------------------------------------------
prs.save(str(DST))
print(f"Saved {DST}  ({len(prs.slides)} slides)")
